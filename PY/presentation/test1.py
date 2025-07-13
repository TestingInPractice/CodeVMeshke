from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Создаем презентацию
prs = Presentation()

# Настройка шрифта и цвета для заголовков и текста
def set_font(run, size=Pt(24), bold=True, color=RGBColor(0, 0, 0)):
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color

def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[5]  # Пустой слайд с заголовком
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    set_font(title_placeholder.text_frame.paragraphs[0].runs[0], size=Pt(32), bold=True)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT

# Данные для слайдов
slides_data = [
    ("Тестирование на практике - базовый курс", [
        "14 модулей, 111 уроков, 815 шагов",
        "7596 учащихся",
        "Рейтинг: 4.5 (88 отзывов)"
    ]),
    ("О курсе", [
        "Обучение азам тестирования с практическими заданиями.",
        "Постепенное усложнение материала.",
        "Реальные задачи из вакансий HHru и рабочих будней тестировщиков."
    ]),
    ("Для кого курс", [
        "Начинающие тестировщики",
        "Аналитики",
        "Разработчики"
    ]),
    ("Чему вы научитесь", [
        "Техники тест-дизайна",
        "Работа с тестовыми артефактами",
        "Использование Devtools и Postman",
        "Основы SQL для работы с базами данных"
    ]),
    ("Программа курса (обзор)", [
        "Профессия тестировщика",
        "Тестовые артефакты",
        "API тестирование",
        "SQL и подготовка к собеседованиям"
    ]),
    ("Отзывы и статистика", [
        "7596 учащихся",
        "Рейтинг 4.5 из 5",
        "88 отзывов"
    ]),
    ("Требования и формат обучения", [
        "Начальный уровень",
        "4 часа в неделю",
        "Умение пользоваться ПК на уровне продвинутого пользователя"
    ]),
    ("Практическая направленность", [
        "Реальные задачи из вакансий и рабочих будней",
        "Обязательные практические задания на каждый шаг"
    ]),
    ("Дополнительные курсы и поддержка", [
        "Подготовка к собеседованиям Junior тестировщика",
        "Расширенный курс по тестированию",
        "Контакты: https://forms.gle/aAYrdyNFg8EmSvSZ7",
        "Telegram: @halyapinvv"
    ]),
    ("Призыв к действию", [
        "Запишитесь на курс сегодня!",
        "Начните свой путь в тестировании с уверенностью."
    ])
]

# Создаем слайды
for title, content in slides_data:
    add_slide(title, content)

# Сохраняем презентацию
prs.save("Testirovanie_na_praktike_Bazovyi_kurs.pptx")

print("Презентация успешно создана: Testirovanie_na_praktike_Bazovyi_kurs.pptx")
